"""Inspect non-tabular resources: PDF, Office docs, HTML/XML, images, ZIP.

Companion to :mod:`file_parser`: while the parser turns tabular files into
DataFrames, the inspector produces an LLM-friendly *description* of everything
else found on agridata.tn (page counts, text previews, image dimensions,
archive listings...). Heavy libraries (pypdf, python-docx, python-pptx,
Pillow) are imported lazily on first use.

Everything is in-memory; nothing touches the filesystem.
"""

from __future__ import annotations

import io
import json
import logging
import re
import zipfile
from typing import Any

import defusedxml.ElementTree as ET

from datagouv_mcp_tn.helpers.file_parser import (
    decode_text_best_effort,
    human_size,
    is_tabular_filename,
)

logger = logging.getLogger(__name__)

TEXT_PREVIEW_CHARS = 600
ARCHIVE_MAX_ENTRIES = 20


# --- magic-byte sniffing for unreliable metadata ('test', 'pb', 'word'...) ---

_MAGIC_SIGNATURES = (
    (b"%PDF", "pdf"),
    (b"PK\x03\x04", "zip"),  # zip, docx, pptx, xlsx
    (b"PK\x05\x06", "zip"),  # empty zip (EOCD only)
    (b"\xd0\xcf\x11\xe0", "ole2"),  # legacy Office (.doc, .xls)
    (b"\x89PNG", "png"),
    (b"GIF8", "gif"),
    (b"\xff\xd8\xff", "jpeg"),
)

_REPLACEMENT_RATIO_LIMIT = 0.02


def sniff_kind(content: bytes) -> str | None:
    """Identify a file family from its magic bytes."""
    for signature, kind in _MAGIC_SIGNATURES:
        if content.startswith(signature):
            return kind
    stripped = content.lstrip()
    if stripped.startswith((b"<", b"{", b"[")):
        return "markup"
    if b"\x00" not in content[:4096]:
        return "text"
    return None


# --- per-family inspectors ---------------------------------------------------


def _inspect_pdf(content: bytes) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    lines = [f"Pages: {len(reader.pages)}"]
    metadata: dict[str, Any] = reader.metadata or {}
    title = str(metadata.get("/Title") or "").strip()
    author = str(metadata.get("/Author") or "").strip()
    if title:
        lines.append(f"Title: {title}")
    if author:
        lines.append(f"Author: {author}")

    text_parts: list[str] = []
    for page in reader.pages:
        try:
            text_parts.append(page.extract_text() or "")
        except Exception:  # nosec B112
            continue
        if sum(len(part) for part in text_parts) >= TEXT_PREVIEW_CHARS:
            break
    preview = " ".join(" ".join(text_parts).split())
    if preview:
        shown = preview[:TEXT_PREVIEW_CHARS]
        suffix = "..." if len(preview) > TEXT_PREVIEW_CHARS else ""
        lines.extend(["", f"Text preview: {shown}{suffix}"])
    else:
        lines.append("No extractable text (scanned image PDF?).")
    return lines


def _inspect_docx(content: bytes) -> list[str]:
    import docx

    document = docx.Document(io.BytesIO(content))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    tables = len(document.tables)
    lines = [
        f"Paragraphs: {len(paragraphs)} · Tables: {tables}",
        f"Words (approx.): {sum(len(p.split()) for p in paragraphs)}",
    ]
    if paragraphs:
        preview = " ".join(paragraphs)
        shown = preview[:TEXT_PREVIEW_CHARS]
        suffix = "..." if len(preview) > TEXT_PREVIEW_CHARS else ""
        lines.extend(["", f"Text preview: {shown}{suffix}"])
    return lines


def _inspect_pptx(content: bytes) -> list[str]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    slides = list(presentation.slides)
    lines = [f"Slides: {len(slides)}"]
    titles: list[str] = []
    for index, slide in enumerate(slides, 1):
        texts = [
            text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and (text := str(getattr(shape, "text", ""))).strip()
        ]
        if texts:
            titles.append(f"{index}. {texts[0][:80]}")
    if titles:
        lines.extend(["", "Slide titles:"])
        shown = titles[:10]
        lines.extend(f"- {title}" for title in shown)
        hidden = len(titles) - len(shown)
        if hidden > 0:
            lines.append(f"... and {hidden} more slide(s) with text")
    return lines


def _inspect_image(content: bytes, fmt: str) -> list[str]:
    from PIL import Image

    with Image.open(io.BytesIO(content)) as image:
        width, height = image.size
        lines = [
            f"Dimensions: {width}×{height}",
            f"Mode: {image.mode} · Frames: {getattr(image, 'n_frames', 1)}",
        ]
        if fmt == "gif":
            animated = getattr(image, "n_frames", 1) > 1
            lines.append(f"Animated: {'yes' if animated else 'no'}")
    lines.append("Images cannot be analyzed further as text; fetch the URL to view it.")
    return lines


def _inspect_svg(content: bytes) -> list[str]:
    root = ET.fromstring(re.sub(rb'xmlns="[^"]+"', b"", content))
    view_box = root.get("viewBox")
    dimensions = (
        root.get("width") and root.get("height") and f"{root.get('width')}×{root.get('height')}"
    )
    shapes = sum(
        len(root.findall(f".//{{*}}{tag}")) for tag in ("path", "rect", "circle", "polygon", "text")
    )
    lines = ["Vector image (SVG)."]
    if dimensions:
        lines.append(f"Declared size: {dimensions}")
    elif view_box:
        lines.append(f"ViewBox: {view_box}")
    lines.append(f"Shape elements: {shapes}")
    return lines


def _inspect_html(text: str) -> list[str]:
    """Structural summary of an HTML page via Scrapy selectors."""
    from scrapy.selector import Selector

    selector = Selector(text=text)
    lines = ["Web page (HTML)."]

    title = selector.css("title::text").get()
    if title and title.strip():
        lines.append(f"Title: {' '.join(title.split())}")

    description = selector.css("meta[name=description]::attr(content)").get()
    if description:
        lines.append(f"Description: {' '.join(description.split())[:200]}")

    counts = {
        "links": len(selector.css("a").getall()),
        "tables": len(selector.css("table").getall()),
        "images": len(selector.css("img").getall()),
        "forms": len(selector.css("form").getall()),
    }
    lines.append(" · ".join(f"{name.capitalize()}: {count}" for name, count in counts.items()))

    headings = [
        h.strip() for h in selector.css("h1 ::text, h2 ::text, h3 ::text").getall() if h.strip()
    ]
    if headings:
        lines.append(f"Headings: {', '.join(headings[:8])}")

    # Drop scripts/styles so the preview only shows visible content.
    selector.css("script, style").drop()
    body_text = (
        selector.xpath("//body//text()").getall() or selector.xpath("//html//text()").getall()
    )
    plain = " ".join(" ".join(body_text).split())
    if plain:
        shown = plain[:TEXT_PREVIEW_CHARS]
        suffix = "..." if len(plain) > TEXT_PREVIEW_CHARS else ""
        lines.extend(["", f"Text preview: {shown}{suffix}"])
    return lines


def _inspect_markup(content: bytes) -> list[str]:
    """HTML / XML / KML / GPX: structural summaries."""
    text = decode_text_best_effort(content).lstrip()
    lowered = text.lower()

    if lowered.startswith("<!doctype html") or lowered.startswith("<html"):
        return _inspect_html(text)

    try:
        root = ET.fromstring(re.sub(rb'xmlns="[^"]+"', b"", content))
    except ET.ParseError:
        return ["Markup file could not be parsed as XML.", "", f"Text preview: {text[:200]}"]

    tag = root.tag.split("}")[-1]
    lines = [f"XML document (root: <{tag}>)."]
    children = list(root)
    lines.append(f"Top-level elements: {len(children)}")
    placemarks = root.findall(".//{*}Placemark") or root.findall(".//Placemark")
    if placemarks:
        lines.append(f"Placemarks: {len(placemarks)} (geographic KML data)")
        names = [p.findtext("{*}name") or p.findtext("name") for p in placemarks[:5]]
        labels = [n for n in names if n]
        if labels:
            lines.append(f"First places: {', '.join(labels)}")
    else:
        counts: dict[str, int] = {}
        for element in root.iter():
            name = element.tag.split("}")[-1]
            counts[name] = counts.get(name, 0) + 1
        top = sorted(counts.items(), key=lambda item: -item[1])[:5]
        if top:
            lines.append("Most common elements: " + ", ".join(f"<{n}>×{c}" for n, c in top))
    return lines


def _inspect_text(content: bytes) -> list[str]:
    """TXT and unknown-but-textual payloads; hint at delimited data."""
    text = decode_text_best_effort(content).lstrip("\ufeff")
    all_lines = text.splitlines()
    first_line = all_lines[0] if all_lines else ""
    delimiters = [d for d in (";", "\t", ",", "|") if first_line.count(d) >= 2]
    lines = [f"Plain text · Lines: {len(all_lines)} · Characters: {len(text)}"]
    if delimiters:
        lines.append(
            f"Delimiters detected on first line ({', '.join(delimiters)}): this may be "
            "CSV-like data saved with a .txt extension."
        )
    shown = text[:TEXT_PREVIEW_CHARS].replace("\n", " ⏎ ")
    lines.extend(["", f"Text preview: {shown}..."])
    return lines


def _inspect_zip(content: bytes) -> list[str]:
    """ZIP (and KMZ): listing + guidance about inner tabular members."""
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        infos = archive.infolist()
        total_uncompressed = sum(info.file_size for info in infos)

        tabular = [info for info in infos if is_tabular_filename(info.filename)]
        lines = [
            f"Archive entries: {len(infos)} · Uncompressed size: {_human(total_uncompressed)}",
        ]
        if tabular:
            names = ", ".join(info.filename for info in tabular[:5])
            more = f" (+{len(tabular) - 5} more)" if len(tabular) > 5 else ""
            lines.append(
                f"Contains {len(tabular)} tabular member(s): {names}{more}. "
                "Members are not individually downloadable via the portal."
            )
        lines.extend(["", "Entries:"])
        shown = infos[:ARCHIVE_MAX_ENTRIES]
        for info in shown:
            lines.append(f"- {info.filename} ({_human(info.file_size)})")
        hidden = len(infos) - len(shown)
        if hidden > 0:
            lines.append(f"... and {hidden} more entr(y/ies)")

        kmz_placemarks = _kmz_placemark_count(content)
        if kmz_placemarks is not None:
            lines.insert(1, f"KMZ archive (KML inside) · Placemarks: {kmz_placemarks}")
        return lines


KMZ_MAX_KML_BYTES = 10 * 1024 * 1024


def _kmz_placemark_count(content: bytes) -> int | None:
    """Count Placemarks in the KML member of a KMZ, with a decompression cap."""
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            info = next(
                (i for i in archive.infolist() if i.filename.lower().endswith(".kml")), None
            )
            if info is None or info.file_size > KMZ_MAX_KML_BYTES:
                return None
            kml_content = archive.read(info)
            root = ET.fromstring(re.sub(rb'xmlns="[^"]+"', b"", kml_content))
            placemarks = root.findall(
                ".//{http://www.opengis.net/kml/2.2}Placemark"
            ) or root.findall(".//Placemark")
            return len(placemarks)
    except Exception:
        return None


def _human(num_bytes: float) -> str:
    return human_size(int(num_bytes))


# --- entry point --------------------------------------------------------------


def inspect_non_tabular(content: bytes, fmt_hint: str) -> tuple[str, list[str]]:
    """Describe any non-tabular payload.

    Args:
        content: Raw downloaded bytes.
        fmt_hint: Normalized format string from the portal ('pdf', 'jpg',
            'word', 'test'... anything the tabular parser rejects).

    Returns:
        ``(kind, lines)`` where ``kind`` identifies the detected family
        (document/image/archive/markup/text/binary).
    """
    if not content:
        return "binary", ["The file is empty (0 bytes)."]

    kind = sniff_kind(content)

    if kind == "zip":
        # DOCX/PPTX/XLSX are ZIP containers; the portal hint disambiguates.
        if fmt_hint in ("docx", "word"):
            return "document", _inspect_docx(content)
        if fmt_hint == "pptx":
            return "document", _inspect_pptx(content)
        return "archive", _inspect_zip(content)
    if kind == "pdf":
        return "document", _inspect_pdf(content)
    if kind == "png":
        return "image", _inspect_image(content, "png")
    if kind == "gif":
        return "image", _inspect_image(content, "gif")
    if kind == "jpeg":
        return "image", _inspect_image(content, "jpg")
    if kind == "ole2":
        return "document", [
            "Legacy Microsoft Office binary (OLE2), most likely Word or Excel.",
            "These files cannot be parsed safely; open them via their URL.",
        ]

    # Not identifiable by magic bytes — trust the portal format hint.
    if kind == "markup":
        if fmt_hint == "svg":
            return "image", _inspect_svg(content)
        if fmt_hint == "json":
            raw_text = decode_text_best_effort(content).strip()
            # Avoid re-serializing huge payloads just for a 400-char preview.
            preview_source = (
                json.dumps(json.loads(raw_text), ensure_ascii=False)
                if len(raw_text) <= 4000
                else raw_text
            )
            return "data", [
                "JSON document (non-tabular).",
                "",
                f"Preview: {preview_source[:400]}...",
            ]
        return "markup", _inspect_markup(content)
    if kind == "text":
        return "text", _inspect_text(content)

    return "binary", [
        f"Unrecognized binary format (portal says '{fmt_hint}', {len(content)} bytes).",
        "Open it directly via its download URL.",
    ]
