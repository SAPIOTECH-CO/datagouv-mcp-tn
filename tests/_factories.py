"""In-test generators producing real PDF/DOCX/PPTX/PNG/ZIP/KMZ fixtures."""

import io
import zipfile
from typing import Any

CSV_BYTES = b"ville,population\nTunis,1056247\nSfax,265131\nSousse,221530\n"


def make_pdf(pages: int = 2) -> bytes:
    from pypdf import PdfWriter

    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=200, height=200)
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def make_docx() -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph("Population de la Tunisie par gouvernorat.")
    document.add_paragraph("Deuxieme paragraphe de donnees.")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_pptx() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title: Any = slide.shapes.title
    assert title is not None
    title.text = "Recensement 2024"
    body: Any = slide.placeholders[1]
    body.text = "Contenu"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_png() -> bytes:
    from PIL import Image

    image = Image.new("RGB", (64, 32), color="red")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def make_zip(members: dict[str, bytes]) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        for name, payload in members.items():
            archive.writestr(name, payload)
    return buffer.getvalue()


def make_kmz(placemarks: int = 3) -> bytes:
    kml = (
        '<?xml version="1.0"?><kml xmlns="http://www.opengis.net/kml/2.2"><Document>'
        + "".join(
            f"<Placemark><name>Gouv {i}</name>"
            "<Point><coordinates>10,36</coordinates></Point></Placemark>"
            for i in range(placemarks)
        )
        + "</Document></kml>"
    )
    return make_zip({"doc.kml": kml.encode()})


HTML_PAGE = b"""<!DOCTYPE html>
<html><head><title>Statistiques</title></head>
<body><p>Budget de l'Etat 2024 en dinars.</p><a href="#">lien</a></body></html>"""

KML_DOC = (
    b'<?xml version="1.0"?><kml xmlns="http://www.opengis.net/kml/2.2">'
    b"<Document><Placemark><name>Tunis</name></Placemark></Document></kml>"
)

SVG_DOC = (
    b'<?xml version="1.0"?><svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 50">'
    b'<path d="M0 0"/><rect width="10" height="10"/><circle r="5"/></svg>'
)
